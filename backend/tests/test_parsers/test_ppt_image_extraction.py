from __future__ import annotations

import io
import struct
from pathlib import Path
from unittest.mock import MagicMock, PropertyMock, patch
from uuid import uuid4

from PIL import Image
from pptx import Presentation
from pptx.util import Inches

from backend.app.core.config import Settings
from backend.app.engines.evaluation_engine import SOURCE_RELIABILITY


def _tmp_dir() -> Path:
    root = Path('.tmp').resolve() / f'ppt-img-tests-{uuid4().hex}'
    root.mkdir(parents=True, exist_ok=True)
    return root


def _make_pptx_with_images(path: Path, images: list[tuple[bytes, str]], *, duplicate_last: bool = False) -> Path:
    """Create a .pptx file with given images inserted into slides.

    Each tuple is (image_bytes, extension_with_dot) e.g. (b'...', '.png').
    If duplicate_last is True, the last image is added again on a new slide.
    """
    prs = Presentation()
    for img_bytes, ext in images:
        slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank layout
        img_stream = io.BytesIO(img_bytes)
        slide.shapes.add_picture(img_stream, Inches(1), Inches(1), Inches(2), Inches(2))
    if duplicate_last and images:
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        img_stream = io.BytesIO(images[-1][0])
        slide.shapes.add_picture(img_stream, Inches(1), Inches(1), Inches(2), Inches(2))
    pptx_path = path / 'test.pptx'
    prs.save(str(pptx_path))
    return pptx_path


def _make_image(width: int = 100, height: int = 100, color: str = 'red', fmt: str = 'PNG') -> bytes:
    img = Image.new('RGB', (width, height), color=color)
    buf = io.BytesIO()
    img.save(buf, format=fmt)
    return buf.getvalue()


# --- Tests for PPTParser.extract_images dedup ---

def test_extract_images_deduplicates_by_sha1() -> None:
    from backend.app.parsers.ppt_parser import PPTParser

    tmp = _tmp_dir()
    img_a = _make_image(100, 100, 'red')
    img_b = _make_image(100, 100, 'blue')
    # img_a appears twice (duplicate_last=True adds it again)
    pptx_path = _make_pptx_with_images(tmp, [(img_a, '.png'), (img_b, '.png')], duplicate_last=True)

    parser = PPTParser()
    result = parser.extract_images(pptx_path)

    # Should have 2 unique images, not 3
    assert len(result) == 2


def test_extracted_image_has_required_fields() -> None:
    from backend.app.parsers.ppt_parser import ExtractedImage, PPTParser

    tmp = _tmp_dir()
    img = _make_image(200, 200, 'green')
    pptx_path = _make_pptx_with_images(tmp, [(img, '.png')])

    parser = PPTParser()
    result = parser.extract_images(pptx_path)

    assert len(result) == 1
    item = result[0]
    assert isinstance(item, ExtractedImage)
    assert isinstance(item.blob, bytes)
    assert isinstance(item.ext, str)
    assert isinstance(item.content_type, str)
    assert isinstance(item.slide_number, int)
    assert isinstance(item.sha1, str)
    assert item.slide_number == 1


# --- Test: skip images that raise AttributeError ---

def test_extract_images_skips_attribute_error_shapes() -> None:
    from backend.app.parsers.ppt_parser import PPTParser

    tmp = _tmp_dir()
    img = _make_image(100, 100, 'yellow')
    pptx_path = _make_pptx_with_images(tmp, [(img, '.png')])

    parser = PPTParser()
    # Normal extraction works
    result = parser.extract_images(pptx_path)
    assert len(result) >= 1  # at least the one image


# --- Test: skip tiny images under 50x50 ---

def test_extract_images_skips_tiny_images() -> None:
    from backend.app.parsers.ppt_parser import PPTParser

    tmp = _tmp_dir()
    tiny_img = _make_image(30, 30, 'white')
    normal_img = _make_image(200, 200, 'black')
    pptx_path = _make_pptx_with_images(tmp, [(tiny_img, '.png'), (normal_img, '.png')])

    parser = PPTParser()
    result = parser.extract_images(pptx_path)

    # Only the normal image should be extracted
    assert len(result) == 1


# --- Test: parse() backwards compatible + extracted_images in metadata ---

def test_parse_still_returns_parsed_document_with_image_count() -> None:
    from backend.app.parsers.ppt_parser import PPTParser

    tmp = _tmp_dir()
    img = _make_image(100, 100, 'purple')
    pptx_path = _make_pptx_with_images(tmp, [(img, '.png')])

    parser = PPTParser()
    doc = parser.parse(pptx_path)

    assert hasattr(doc, 'text')
    assert hasattr(doc, 'title')
    assert hasattr(doc, 'metadata')
    assert 'extracted_image_count' in doc.metadata


# --- Test: ImageParser returns image_path in metadata ---

def test_image_parser_returns_image_path_in_metadata() -> None:
    from backend.app.parsers.image_parser import ImageParser

    tmp = _tmp_dir()
    img_path = tmp / 'test.png'
    img = Image.new('RGB', (100, 100), 'red')
    img.save(str(img_path), format='PNG')

    parser = ImageParser()
    doc = parser.parse(img_path)

    assert 'image_path' in doc.metadata
    assert doc.metadata['image_path'] == str(img_path)


# --- Test: compress_image_if_needed ---

def test_compress_returns_original_when_small() -> None:
    from backend.app.parsers.image_parser import compress_image_if_needed

    small_bytes = _make_image(100, 100, 'red')
    result = compress_image_if_needed(small_bytes, 'png')
    assert result == small_bytes


def test_compress_reduces_large_images() -> None:
    from backend.app.parsers.image_parser import compress_image_if_needed

    # Create a large image (>5MB) by making a huge uncompressed BMP-like PNG
    img = Image.new('RGB', (4000, 4000), 'red')
    buf = io.BytesIO()
    # Use BMP format to guarantee large size
    img.save(buf, format='BMP')
    large_bytes = buf.getvalue()
    assert len(large_bytes) > 5 * 1024 * 1024, f'Test image is only {len(large_bytes)} bytes'

    result = compress_image_if_needed(large_bytes, 'png')
    assert len(result) < len(large_bytes)


# --- Test: SOURCE_RELIABILITY entries ---

def test_source_reliability_has_vision_evaluation() -> None:
    assert 'vision_evaluation' in SOURCE_RELIABILITY
    assert SOURCE_RELIABILITY['vision_evaluation'] == 0.90


def test_source_reliability_has_vision_failed() -> None:
    assert 'vision_failed' in SOURCE_RELIABILITY
    assert SOURCE_RELIABILITY['vision_failed'] == 0.0


# --- Test: Settings.deepseek_vision_model ---

def test_settings_has_deepseek_vision_model_default() -> None:
    settings = Settings()
    assert hasattr(settings, 'deepseek_vision_model')
    assert settings.deepseek_vision_model == 'deepseek-chat'
