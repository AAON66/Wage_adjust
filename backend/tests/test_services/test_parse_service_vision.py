from __future__ import annotations

import io
from dataclasses import dataclass
from pathlib import Path
from unittest.mock import MagicMock, patch
from uuid import uuid4

from PIL import Image
from pptx import Presentation
from pptx.util import Inches

from backend.app.core.config import Settings
from backend.app.core.database import create_db_engine, create_session_factory, init_database
from backend.app.models import load_model_modules
from backend.app.models.employee import Employee
from backend.app.models.evaluation_cycle import EvaluationCycle
from backend.app.models.submission import EmployeeSubmission
from backend.app.models.uploaded_file import UploadedFile
from backend.app.services.llm_service import DeepSeekCallResult
from backend.app.services.parse_service import ParseService


def _build_context() -> tuple[Settings, object]:
    temp_root = Path('.tmp').resolve()
    temp_root.mkdir(parents=True, exist_ok=True)
    database_path = (temp_root / f'parse-vision-{uuid4().hex}.db').as_posix()
    uploads_path = (temp_root / f'uploads-vision-{uuid4().hex}').as_posix()
    settings = Settings(
        database_url=f'sqlite+pysqlite:///{database_path}',
        storage_base_dir=uploads_path,
        deepseek_api_key='test_key',
    )
    settings.deepseek_require_real_call_for_parsing = False
    load_model_modules()
    engine = create_db_engine(settings)
    init_database(engine)
    return settings, create_session_factory(settings)


def _seed_submission(settings, session_factory):
    db = session_factory()
    employee = Employee(
        employee_no=f'EMP-V-{uuid4().hex[:6]}',
        name='Vision User',
        department='Engineering',
        job_family='Platform',
        job_level='P5',
        status='active',
    )
    cycle = EvaluationCycle(name='2026 Vision', review_period='2026', budget_amount='1000.00', status='draft')
    db.add_all([employee, cycle])
    db.commit()
    db.refresh(employee)
    db.refresh(cycle)

    submission = EmployeeSubmission(employee_id=employee.id, cycle_id=cycle.id, status='submitted')
    db.add(submission)
    db.commit()
    db.refresh(submission)
    return db, submission


def _create_test_pptx_with_images(target_path: Path, num_images: int = 2) -> None:
    """Create a real .pptx file with embedded images."""
    prs = Presentation()
    for i in range(num_images):
        slide = prs.slides.add_slide(prs.slide_layouts[5])  # blank layout
        # Create a unique image for each slide
        img_buf = io.BytesIO()
        img = Image.new('RGB', (200, 200), color=(100 + i * 50, 50, 50))
        img.save(img_buf, format='PNG')
        img_buf.seek(0)
        slide.shapes.add_picture(img_buf, Inches(1), Inches(1), Inches(3), Inches(3))
        # Add some text too
        from pptx.util import Pt
        txBox = slide.shapes.add_textbox(Inches(1), Inches(5), Inches(5), Inches(1))
        txBox.text_frame.text = f'Slide {i + 1}: AI workflow automation results'
    prs.save(str(target_path))


def _create_test_png(target_path: Path) -> None:
    """Create a real .png image file."""
    img = Image.new('RGB', (200, 200), color=(100, 150, 200))
    img.save(str(target_path), format='PNG')


def _mock_deepseek_success(image_bytes, ext, mime_type, *, context=None):
    """Return a successful vision evaluation result."""
    return DeepSeekCallResult(
        payload={
            'description': 'Screenshot showing AI workflow dashboard with metrics',
            'quality_score': 4,
            'dimension_relevance': {'ai_tool_mastery': 0.8, 'ai_result_conversion': 0.6},
        },
        used_fallback=False,
        provider='deepseek',
    )


def _mock_deepseek_fallback(image_bytes, ext, mime_type, *, context=None):
    """Return a fallback result (vision API unavailable)."""
    return DeepSeekCallResult(
        payload={'description': 'Vision unavailable', 'quality_score': 0, 'dimension_relevance': {}},
        used_fallback=True,
        provider='deepseek',
        reason='vision_api_unavailable',
    )


def _seed_with_pptx(num_images: int = 2):
    settings, session_factory = _build_context()
    db, submission = _seed_submission(settings, session_factory)

    storage_dir = Path(settings.storage_base_dir).resolve()
    storage_dir.mkdir(parents=True, exist_ok=True)
    submission_dir = storage_dir / submission.id
    submission_dir.mkdir(parents=True, exist_ok=True)
    target_file = submission_dir / 'presentation.pptx'
    _create_test_pptx_with_images(target_file, num_images=num_images)

    uploaded_file = UploadedFile(
        submission_id=submission.id,
        file_name='presentation.pptx',
        file_type='pptx',
        storage_key=f'{submission.id}/presentation.pptx',
        parse_status='pending',
    )
    db.add(uploaded_file)
    db.commit()
    db.refresh(uploaded_file)
    return settings, db, submission, uploaded_file


def _seed_with_png():
    settings, session_factory = _build_context()
    db, submission = _seed_submission(settings, session_factory)

    storage_dir = Path(settings.storage_base_dir).resolve()
    storage_dir.mkdir(parents=True, exist_ok=True)
    submission_dir = storage_dir / submission.id
    submission_dir.mkdir(parents=True, exist_ok=True)
    target_file = submission_dir / 'screenshot.png'
    _create_test_png(target_file)

    uploaded_file = UploadedFile(
        submission_id=submission.id,
        file_name='screenshot.png',
        file_type='png',
        storage_key=f'{submission.id}/screenshot.png',
        parse_status='pending',
    )
    db.add(uploaded_file)
    db.commit()
    db.refresh(uploaded_file)
    return settings, db, submission, uploaded_file


def _seed_with_txt():
    settings, session_factory = _build_context()
    db, submission = _seed_submission(settings, session_factory)

    storage_dir = Path(settings.storage_base_dir).resolve()
    storage_dir.mkdir(parents=True, exist_ok=True)
    submission_dir = storage_dir / submission.id
    submission_dir.mkdir(parents=True, exist_ok=True)
    target_file = submission_dir / 'notes.txt'
    target_file.write_text('AI workflow improvements documentation.', encoding='utf-8')

    uploaded_file = UploadedFile(
        submission_id=submission.id,
        file_name='notes.txt',
        file_type='txt',
        storage_key=f'{submission.id}/notes.txt',
        parse_status='pending',
    )
    db.add(uploaded_file)
    db.commit()
    db.refresh(uploaded_file)
    return settings, db, submission, uploaded_file


def test_pptx_vision_creates_evidence_items_for_each_image() -> None:
    """parse_file() for a .pptx with 2 images creates 2 vision EvidenceItems + text EvidenceItems."""
    settings, db, submission, uploaded_file = _seed_with_pptx(num_images=2)
    try:
        mock_deepseek = MagicMock()
        mock_deepseek.evaluate_image_vision = MagicMock(side_effect=_mock_deepseek_success)
        service = ParseService(db, settings, deepseek_service=mock_deepseek)
        parsed_file, evidence_count = service.parse_file(uploaded_file)

        db.refresh(submission)
        assert parsed_file.parse_status == 'parsed'
        vision_items = [e for e in submission.evidence_items if e.source_type == 'vision_evaluation']
        assert len(vision_items) == 2
        # Also has text evidence
        text_items = [e for e in submission.evidence_items if e.source_type != 'vision_evaluation']
        assert len(text_items) >= 1
        assert evidence_count == len(submission.evidence_items)
    finally:
        db.close()


def test_vision_evidence_has_correct_source_type() -> None:
    """Each vision EvidenceItem has source_type='vision_evaluation'."""
    settings, db, submission, uploaded_file = _seed_with_pptx(num_images=1)
    try:
        mock_deepseek = MagicMock()
        mock_deepseek.evaluate_image_vision = MagicMock(side_effect=_mock_deepseek_success)
        service = ParseService(db, settings, deepseek_service=mock_deepseek)
        service.parse_file(uploaded_file)

        db.refresh(submission)
        vision_items = [e for e in submission.evidence_items if e.source_type == 'vision_evaluation']
        assert len(vision_items) == 1
        assert vision_items[0].source_type == 'vision_evaluation'
    finally:
        db.close()


def test_vision_evidence_has_correct_confidence_score() -> None:
    """Each vision EvidenceItem has confidence_score = quality_score / 5.0."""
    settings, db, submission, uploaded_file = _seed_with_pptx(num_images=1)
    try:
        mock_deepseek = MagicMock()
        mock_deepseek.evaluate_image_vision = MagicMock(side_effect=_mock_deepseek_success)
        service = ParseService(db, settings, deepseek_service=mock_deepseek)
        service.parse_file(uploaded_file)

        db.refresh(submission)
        vision_items = [e for e in submission.evidence_items if e.source_type == 'vision_evaluation']
        assert len(vision_items) == 1
        # quality_score=4, so confidence = 4/5.0 = 0.8
        assert vision_items[0].confidence_score == 0.8
    finally:
        db.close()


def test_vision_evidence_metadata_for_ppt_embedded() -> None:
    """Vision EvidenceItem metadata contains expected fields for PPT-embedded images."""
    settings, db, submission, uploaded_file = _seed_with_pptx(num_images=1)
    try:
        mock_deepseek = MagicMock()
        mock_deepseek.evaluate_image_vision = MagicMock(side_effect=_mock_deepseek_success)
        service = ParseService(db, settings, deepseek_service=mock_deepseek)
        service.parse_file(uploaded_file)

        db.refresh(submission)
        vision_items = [e for e in submission.evidence_items if e.source_type == 'vision_evaluation']
        assert len(vision_items) == 1
        meta = vision_items[0].metadata_json
        assert 'file_id' in meta
        assert 'storage_key' in meta
        assert meta['vision_quality_score'] == 4
        assert 'vision_description' in meta
        assert isinstance(meta['vision_dimension_relevance'], dict)
        assert meta['slide_number'] == 1
        assert meta['image_source'] == 'ppt_embedded'
    finally:
        db.close()


def test_standalone_png_creates_vision_evidence() -> None:
    """parse_file() for a standalone .png creates 1 vision EvidenceItem with image_source='standalone_upload'."""
    settings, db, submission, uploaded_file = _seed_with_png()
    try:
        mock_deepseek = MagicMock()
        mock_deepseek.evaluate_image_vision = MagicMock(side_effect=_mock_deepseek_success)
        mock_deepseek.extract_image_text = MagicMock(return_value=DeepSeekCallResult(
            payload={'has_text': False, 'extracted_text': ''},
            used_fallback=True,
            provider='deepseek',
        ))
        service = ParseService(db, settings, deepseek_service=mock_deepseek)
        parsed_file, evidence_count = service.parse_file(uploaded_file)

        db.refresh(submission)
        assert parsed_file.parse_status == 'parsed'
        vision_items = [e for e in submission.evidence_items if e.source_type == 'vision_evaluation']
        assert len(vision_items) == 1
        meta = vision_items[0].metadata_json
        assert meta['image_source'] == 'standalone_upload'
    finally:
        db.close()


def test_vision_failure_isolation_one_fails_other_succeeds() -> None:
    """When vision API fails for one image, that image gets vision_failed, other images still succeed (D-07)."""
    settings, db, submission, uploaded_file = _seed_with_pptx(num_images=2)
    try:
        call_count = 0

        def _side_effect(image_bytes, ext, mime_type, *, context=None):
            nonlocal call_count
            call_count += 1
            if call_count == 1:
                raise RuntimeError('Vision API error')
            return _mock_deepseek_success(image_bytes, ext, mime_type, context=context)

        mock_deepseek = MagicMock()
        mock_deepseek.evaluate_image_vision = MagicMock(side_effect=_side_effect)
        service = ParseService(db, settings, deepseek_service=mock_deepseek)
        service.parse_file(uploaded_file)

        db.refresh(submission)
        failed_items = [e for e in submission.evidence_items if e.source_type == 'vision_failed']
        success_items = [e for e in submission.evidence_items if e.source_type == 'vision_evaluation']
        assert len(failed_items) == 1
        assert len(success_items) == 1
        assert failed_items[0].confidence_score == 0.0
        assert failed_items[0].metadata_json.get('vision_failed') is True
    finally:
        db.close()


def test_vision_used_fallback_creates_vision_failed() -> None:
    """When vision API returns used_fallback=True, the evidence has source_type='vision_failed'."""
    settings, db, submission, uploaded_file = _seed_with_pptx(num_images=1)
    try:
        mock_deepseek = MagicMock()
        mock_deepseek.evaluate_image_vision = MagicMock(side_effect=_mock_deepseek_fallback)
        service = ParseService(db, settings, deepseek_service=mock_deepseek)
        service.parse_file(uploaded_file)

        db.refresh(submission)
        vision_items = [e for e in submission.evidence_items
                        if e.source_type in ('vision_evaluation', 'vision_failed')]
        assert len(vision_items) == 1
        assert vision_items[0].source_type == 'vision_failed'
        assert vision_items[0].confidence_score == 0.0
    finally:
        db.close()


def test_no_deepseek_service_means_no_vision_evidence() -> None:
    """When deepseek_service is None, no vision evaluation happens (graceful skip)."""
    settings, db, submission, uploaded_file = _seed_with_pptx(num_images=2)
    try:
        service = ParseService(db, settings, deepseek_service=None)
        parsed_file, evidence_count = service.parse_file(uploaded_file)

        db.refresh(submission)
        vision_items = [e for e in submission.evidence_items
                        if e.source_type in ('vision_evaluation', 'vision_failed')]
        assert len(vision_items) == 0
        assert parsed_file.parse_status == 'parsed'
    finally:
        db.close()


def test_txt_file_produces_no_vision_evidence() -> None:
    """parse_file() for a .txt file produces no vision evidence (unchanged behavior)."""
    settings, db, submission, uploaded_file = _seed_with_txt()
    try:
        mock_deepseek = MagicMock()
        mock_deepseek.evaluate_image_vision = MagicMock(side_effect=_mock_deepseek_success)
        service = ParseService(db, settings, deepseek_service=mock_deepseek)
        parsed_file, evidence_count = service.parse_file(uploaded_file)

        db.refresh(submission)
        vision_items = [e for e in submission.evidence_items
                        if e.source_type in ('vision_evaluation', 'vision_failed')]
        assert len(vision_items) == 0
        mock_deepseek.evaluate_image_vision.assert_not_called()
    finally:
        db.close()


def test_parse_submission_files_mixed_types() -> None:
    """parse_submission_files() with mixed file types correctly processes each type."""
    # This test creates separate contexts, so just verify individual parsing works
    # for each type through the same service
    settings, db, submission, uploaded_file_pptx = _seed_with_pptx(num_images=1)
    try:
        mock_deepseek = MagicMock()
        mock_deepseek.evaluate_image_vision = MagicMock(side_effect=_mock_deepseek_success)
        service = ParseService(db, settings, deepseek_service=mock_deepseek)

        # Create a txt file in the same submission
        storage_dir = Path(settings.storage_base_dir).resolve()
        submission_dir = storage_dir / submission.id
        txt_path = submission_dir / 'notes.txt'
        txt_path.write_text('AI workflow docs.', encoding='utf-8')
        uploaded_file_txt = UploadedFile(
            submission_id=submission.id,
            file_name='notes.txt',
            file_type='txt',
            storage_key=f'{submission.id}/notes.txt',
            parse_status='pending',
        )
        db.add(uploaded_file_txt)
        db.commit()
        db.refresh(uploaded_file_txt)

        updated_files, total_evidence = service.parse_submission_files([uploaded_file_pptx, uploaded_file_txt])

        db.refresh(submission)
        assert len(updated_files) == 2
        assert all(f.parse_status == 'parsed' for f in updated_files)
        # pptx should have vision evidence, txt should not
        vision_items = [e for e in submission.evidence_items
                        if e.source_type in ('vision_evaluation', 'vision_failed')]
        assert len(vision_items) == 1  # 1 image in pptx
    finally:
        db.close()
