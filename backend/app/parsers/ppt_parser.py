from __future__ import annotations

import io
import logging
from dataclasses import dataclass
from pathlib import Path

from PIL import Image
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

from backend.app.parsers.base_parser import BaseParser, ParsedDocument

logger = logging.getLogger(__name__)

MIN_IMAGE_DIMENSION = 50  # skip icons/bullets


@dataclass
class ExtractedImage:
    blob: bytes
    ext: str
    content_type: str
    slide_number: int
    shape_id: int
    sha1: str


class PPTParser(BaseParser):
    supported_extensions = ('.pptx',)

    def extract_images(self, path: Path) -> list[ExtractedImage]:
        presentation = Presentation(str(path))
        seen_hashes: set[str] = set()
        images: list[ExtractedImage] = []
        for slide_idx, slide in enumerate(presentation.slides, start=1):
            for shape in slide.shapes:
                try:
                    # python-pptx raises NotImplementedError for some shape types
                    # (see python-pptx #929 and similar); skip them gracefully.
                    if shape.shape_type != MSO_SHAPE_TYPE.PICTURE:
                        continue
                except (NotImplementedError, ValueError) as exc:
                    logger.warning('Skipping unrecognized shape on slide %d: %s', slide_idx, exc)
                    continue
                try:
                    image = shape.image
                    if image.sha1 in seen_hashes:
                        continue
                    # Skip tiny images (icons, bullets)
                    try:
                        with Image.open(io.BytesIO(image.blob)) as img:
                            w, h = img.size
                            if w < MIN_IMAGE_DIMENSION or h < MIN_IMAGE_DIMENSION:
                                continue
                    except Exception:
                        pass  # If we can't read dimensions, keep the image
                    seen_hashes.add(image.sha1)
                    images.append(ExtractedImage(
                        blob=image.blob,
                        ext=image.ext,
                        content_type=image.content_type,
                        slide_number=slide_idx,
                        shape_id=shape.shape_id,
                        sha1=image.sha1,
                    ))
                except (AttributeError, ValueError) as exc:
                    logger.warning('Failed to extract image from slide %d shape %s: %s', slide_idx, getattr(shape, 'shape_id', '?'), exc)
        return images

    def parse(self, path: Path) -> ParsedDocument:
        presentation = Presentation(str(path))
        chunks: list[str] = []
        for slide in presentation.slides:
            for shape in slide.shapes:
                if hasattr(shape, 'text') and shape.text:
                    chunks.append(shape.text)

        text = '\n'.join(chunk.strip() for chunk in chunks if chunk.strip())
        # Extract images for metadata count (per D-05)
        extracted_images = self.extract_images(path)
        return ParsedDocument(
            text=text or f'No slide text extracted from {path.name}.',
            title=path.name,
            metadata={
                'slides': len(presentation.slides),
                'extension': path.suffix.lower(),
                'extracted_image_count': len(extracted_images),
            },
        )
